"""Exporte les livrables Markdown en fichiers DOCX et PPTX."""

from __future__ import annotations

import sys
from pathlib import Path

ROOT_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT_DIR))

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt


DOCS_DIR = ROOT_DIR / "docs"
REPORT_MD = DOCS_DIR / "rapport_technique.md"
PRESENTATION_MD = DOCS_DIR / "presentation_poc_rag.md"
REPORT_DOCX = DOCS_DIR / "rapport_technique.docx"
PRESENTATION_PPTX = DOCS_DIR / "presentation_poc_rag.pptx"


def export_report() -> None:
    document = Document()
    for raw_line in REPORT_MD.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("# "):
            document.add_heading(line[2:], level=0)
        elif line.startswith("## "):
            document.add_heading(line[3:], level=1)
        elif line.startswith("```"):
            continue
        elif line[0:2] in {"1.", "2.", "3.", "4.", "5.", "6.", "7.", "8.", "9."}:
            document.add_paragraph(line[3:], style="List Number")
        else:
            document.add_paragraph(line)
    document.save(REPORT_DOCX)


def parse_slides() -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_body: list[str] = []

    for raw_line in PRESENTATION_MD.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if line.startswith("# "):
            continue
        if line.startswith("## "):
            if current_title:
                slides.append((current_title, current_body))
            current_title = line[3:]
            current_body = []
        elif line and current_title:
            current_body.append(line)

    if current_title:
        slides.append((current_title, current_body))
    return slides


def export_presentation() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for title, body_lines in parse_slides():
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()
        for index, line in enumerate(body_lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(24 if index == 0 else 20)
            paragraph.level = 0

    presentation.save(PRESENTATION_PPTX)


def main() -> None:
    export_report()
    export_presentation()
    print(f"Rapport Word enregistré: {REPORT_DOCX}")
    print(f"Présentation PowerPoint enregistrée: {PRESENTATION_PPTX}")


if __name__ == "__main__":
    main()
